import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Palette
    DARK_BG = RGBColor(40, 48, 36)        # Deep forest (#283024)
    LIGHT_BG = RGBColor(248, 246, 240)    # Rice paper cream (#F8F6F0)
    CARD_BG = RGBColor(255, 255, 255)     # Crisp white
    DARK_CARD_BG = RGBColor(52, 62, 47)   # Dark Forest Card
    PRIMARY = RGBColor(93, 112, 82)       # Moss green (#5D7052)
    SECONDARY = RGBColor(193, 140, 93)    # Terracotta (#C18C5D)
    TEXT_DARK = RGBColor(44, 44, 36)      # Deep Loam (#2C2C24)
    TEXT_LIGHT = RGBColor(248, 248, 242)  # Rice Paper (#FDFCF8)
    TEXT_MUTED = RGBColor(135, 135, 125)  # Muted grey
    BORDER_COLOR = RGBColor(225, 220, 210)

    def set_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, category, title, dark=False):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = SECONDARY if dark else PRIMARY

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.75))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(23)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_LIGHT if dark else TEXT_DARK

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    def create_textbox(slide, left, top, width, height):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        return tf

    # ==========================================
    # SLIDE 1: Title Slide (Dark)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, DARK_BG)

    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(2.2), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_CARD_BG
    badge.line.color.rgb = PRIMARY
    badge.text = "🌿 LEARNAI PLATFORM"
    for p in badge.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

    tf = create_textbox(s1, Inches(0.8), Inches(1.9), Inches(11.5), Inches(2.0))
    p = tf.paragraphs[0]
    p.text = "An AI-Powered Personalized\nLearning & Skill-Gap Platform"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT

    tf_sub = create_textbox(s1, Inches(0.8), Inches(4.2), Inches(10.5), Inches(1.2))
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Turning free-text career goals, resume uploads, and skill assessments into structured, explainable, and adaptive learning paths — built on Supabase, pgvector, and Groq LLMs."
    p_sub.font.size = Pt(15.5)
    p_sub.font.color.rgb = RGBColor(200, 205, 195)

    tf_auth = create_textbox(s1, Inches(0.8), Inches(5.7), Inches(6.0), Inches(0.5))
    p_auth = tf_auth.paragraphs[0]
    p_auth.text = "By Code Catalyst"
    p_auth.font.size = Pt(18)
    p_auth.font.bold = True
    p_auth.font.color.rgb = SECONDARY

    tf_foot = create_textbox(s1, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.4))
    p_foot = tf_foot.paragraphs[0]
    p_foot.text = "PROJECT OVERVIEW  ·  AUGUST 2026"
    p_foot.font.size = Pt(11)
    p_foot.font.bold = True
    p_foot.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: The Problem (Light)
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, LIGHT_BG)
    add_header(s2, "The Problem", "Generic course catalogs don't know what you already know")

    col_w = Inches(3.64)
    gap = Inches(0.38)
    top_pos = Inches(1.9)
    card_h = Inches(4.8)

    problems = [
        ("1", "One-Size-Fits-All Catalogs", "Learners scroll through long lists of courses with no sense of pedagogical order, prerequisite validation, or individualized starting points based on existing skills."),
        ("2", "Goals Stay Unstructured", "\"I want to become an ML engineer\" or \"I want to be a Doctor\" never gets translated into concrete, measurable, and verified skill competencies."),
        ("3", "No Visible Progress Logic", "Recommendations feel like a black box — learners cannot see why a course was suggested, how it fits their roadmap, or what tangible gap it closes.")
    ]

    for i, (num, headline, desc) in enumerate(problems):
        left_pos = Inches(0.8) + i * (col_w + gap)
        add_card(s2, left_pos, top_pos, col_w, card_h)

        b_num = s2.shapes.add_shape(MSO_SHAPE.OVAL, left_pos + Inches(0.35), top_pos + Inches(0.4), Inches(0.65), Inches(0.65))
        b_num.fill.solid()
        b_num.fill.fore_color.rgb = PRIMARY
        b_num.line.fill.background()
        b_num.text = num
        b_num.text_frame.paragraphs[0].font.size = Pt(14)
        b_num.text_frame.paragraphs[0].font.bold = True
        b_num.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
        b_num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        tf = create_textbox(s2, left_pos + Inches(0.35), top_pos + Inches(1.35), col_w - Inches(0.7), Inches(3.0))
        p1 = tf.paragraphs[0]
        p1.text = headline
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_DARK

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12.5)
        p2.font.color.rgb = RGBColor(100, 100, 95)
        p2.space_before = Pt(14)

    # ==========================================
    # SLIDE 3: The Solution (Dark)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, DARK_BG)
    add_header(s3, "The Solution", "LearnAI turns goals & resumes into guided, explainable paths", dark=True)

    steps = [
        ("Describe / Upload", "Free-text career goals or resume uploads (.pdf / .docx) across CS & Non-CS domains."),
        ("AI Parses Skills", "Groq LLM extracts canonical skills & benchmarks against 100+ roles in role_skills.json."),
        ("Path Is Generated", "Topological sequencing + pgvector cosine similarity + skill-gap scoring builds the order."),
        ("Track & Adapt", "Interactive verified quizzes, adaptive progress boosting, and a scoped AI learning guide.")
    ]

    col4_w = Inches(2.7)
    gap4 = Inches(0.3)
    top4 = Inches(2.0)
    h4 = Inches(4.6)

    for i, (title, desc) in enumerate(steps):
        left_pos = Inches(0.8) + i * (col4_w + gap4)
        add_card(s3, left_pos, top4, col4_w, h4, bg_color=DARK_CARD_BG, border_color=PRIMARY)

        tag = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos + Inches(0.3), top4 + Inches(0.35), Inches(1.2), Inches(0.35))
        tag.fill.solid()
        tag.fill.fore_color.rgb = SECONDARY
        tag.line.fill.background()
        tag.text = f"STEP {i+1}"
        tag.text_frame.paragraphs[0].font.size = Pt(10)
        tag.text_frame.paragraphs[0].font.bold = True
        tag.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
        tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        tf = create_textbox(s3, left_pos + Inches(0.3), top4 + Inches(1.0), col4_w - Inches(0.6), Inches(3.2))
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15.5)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_LIGHT

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(200, 205, 195)
        p2.space_before = Pt(10)

    # ==========================================
    # SLIDE 4: Architecture (Light - Matched to Screenshot)
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, LIGHT_BG)
    add_header(s4, "Architecture", "A lean, serverless stack")

    # Card 1: Browser
    add_card(s4, Inches(0.8), Inches(1.8), Inches(11.73), Inches(1.15))
    tf1 = create_textbox(s4, Inches(1.05), Inches(1.88), Inches(11.2), Inches(0.95))
    p = tf1.paragraphs[0]
    p.text = "💻 Browser & Client Layer"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p2 = tf1.add_paragraph()
    p2.text = "React 19 + Vite 8 + Tailwind CSS 4 • React Router 7 • TanStack Query • Recharts • PDF.js & Mammoth client-side parsing"
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = RGBColor(100, 100, 95)
    p2.space_before = Pt(3)

    # Card 2: Auth & Postgres (Left)
    add_card(s4, Inches(0.8), Inches(3.15), Inches(5.7), Inches(1.85))
    tf2 = create_textbox(s4, Inches(1.05), Inches(3.25), Inches(5.2), Inches(1.65))
    p = tf2.paragraphs[0]
    p.text = "🔒 Auth & Postgres (RLS)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p2 = tf2.add_paragraph()
    p2.text = "Row-level security on every learner table, pgvector extension for 384-dim course embeddings, auth.users-backed profiles."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = RGBColor(100, 100, 95)
    p2.space_before = Pt(6)

    # Card 3: Edge Functions (Right)
    add_card(s4, Inches(6.83), Inches(3.15), Inches(5.7), Inches(1.85))
    tf3 = create_textbox(s4, Inches(7.08), Inches(3.25), Inches(5.2), Inches(1.65))
    p = tf3.paragraphs[0]
    p.text = "⚡ Edge Functions (Deno/TS)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p2 = tf3.add_paragraph()
    p2.text = "Eight serverless functions handle parsing, vector scoring, path assembly, quiz verification, live progress rescoring, and chat — invoked via supabase.functions.invoke()."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = RGBColor(100, 100, 95)
    p2.space_before = Pt(6)

    # Card 4: Groq LLM
    add_card(s4, Inches(0.8), Inches(5.2), Inches(11.73), Inches(1.2))
    tf4 = create_textbox(s4, Inches(1.05), Inches(5.28), Inches(11.2), Inches(1.0))
    p = tf4.paragraphs[0]
    p.text = "🤖 Groq LLM Inference"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    p2 = tf4.add_paragraph()
    p2.text = "Goal extraction, milestone rationales, verified quiz generation, and scoped AI chat all run through Groq for sub-second low-latency inference."
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = RGBColor(100, 100, 95)
    p2.space_before = Pt(3)

    # Footer note
    tf_f4 = create_textbox(s4, Inches(0.8), Inches(6.6), Inches(11.73), Inches(0.35))
    p = tf_f4.paragraphs[0]
    p.text = "Deployed on Vercel (frontend) + Supabase (backend) — no dedicated servers to manage."
    p.font.size = Pt(10.5)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 5: Core Capabilities (Light)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, LIGHT_BG)
    add_header(s5, "Core Features", "Six capabilities, one continuous experience")

    features = [
        ("Conversational Goal Parsing", "Free-text goals mapped to a canonical skill taxonomy via Groq LLM extraction without hallucination."),
        ("Resume Skill Gap Analyzer", "Upload .pdf / .docx resumes to benchmark qualifications against 100+ roles in role_skills.json."),
        ("Personalized Learning Paths", "Topological curriculum curation combining pgvector similarity with marginal skill-gap scoring."),
        ("Interactive MCQ Verification", "Dynamic 3-question quizzes generated per skill via Groq with server-side grading and security."),
        ("Adaptive Progress & Edit Plan", "Real-time skill proficiency boosting upon completion + full flexibility to prune path items."),
        ("AI Learning Guide", "A scoped mentor that reasons directly over each learner's chat history, verified skills, and roadmap.")
    ]

    card_w5 = Inches(3.64)
    card_h5 = Inches(2.2)
    for idx, (ft_title, ft_desc) in enumerate(features):
        row = idx // 3
        col = idx % 3
        c_left = Inches(0.8) + col * (card_w5 + Inches(0.38))
        c_top = Inches(1.9) + row * (card_h5 + Inches(0.3))

        add_card(s5, c_left, c_top, card_w5, card_h5)
        tf = create_textbox(s5, c_left + Inches(0.25), c_top + Inches(0.2), card_w5 - Inches(0.5), card_h5 - Inches(0.4))
        p = tf.paragraphs[0]
        p.text = ft_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p2 = tf.add_paragraph()
        p2.text = ft_desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = RGBColor(100, 100, 95)
        p2.space_before = Pt(6)

    # ==========================================
    # SLIDE 6: The AI Layer (Dark)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, DARK_BG)
    add_header(s6, "AI Layer", "Eight Edge Functions power every AI moment", dark=True)

    functions_list = [
        ("parse-goal", "Extracts skills & experience level from free-text goals via Groq LLM extraction."),
        ("get-recommendations", "Vector similarity search over course embeddings, blended with skill-gap scoring."),
        ("generate-path", "Greedy topological builder that orders courses into milestone-labeled paths."),
        ("update-progress", "Live boosts skill proficiency on completion and dynamically re-scores remaining items."),
        ("generate-quiz", "Generates 3-question MCQ skill verification quizzes server-side via Groq."),
        ("submit-quiz", "Evaluates submissions, logs attempts to quiz_attempts, and verifies skills in DB."),
        ("ask-assistant", "Context-aware chat assistant scoped to the learner's live profile, skills, and path."),
        ("explain-path-item", "Generates a short, personalized rationale for each recommended milestone course.")
    ]

    fn_w = Inches(5.6)
    fn_h = Inches(1.05)
    for idx, (fn_name, fn_desc) in enumerate(functions_list):
        row = idx % 4
        col = idx // 4
        f_left = Inches(0.8) + col * (fn_w + Inches(0.53))
        f_top = Inches(1.9) + row * (fn_h + Inches(0.25))

        add_card(s6, f_left, f_top, fn_w, fn_h, bg_color=DARK_CARD_BG, border_color=PRIMARY)
        tf = create_textbox(s6, f_left + Inches(0.2), f_top + Inches(0.12), fn_w - Inches(0.4), fn_h - Inches(0.24))
        p = tf.paragraphs[0]
        p.text = f"⚡ {fn_name}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = SECONDARY
        p2 = tf.add_paragraph()
        p2.text = fn_desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = RGBColor(210, 215, 205)
        p2.space_before = Pt(3)

    # ==========================================
    # SLIDE 7: Scoring Engine (Light)
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7, LIGHT_BG)
    add_header(s7, "How Scoring Works", "Every course gets a final, explainable score")

    # Banner
    add_card(s7, Inches(0.8), Inches(1.8), Inches(11.73), Inches(1.3), bg_color=RGBColor(235, 240, 230), border_color=PRIMARY)
    tf_f = create_textbox(s7, Inches(1.0), Inches(1.95), Inches(11.3), Inches(1.0))
    p = tf_f.paragraphs[0]
    p.text = "final_score  =  0.6 × similarity  +  0.4 × gap_score"
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER
    p2 = tf_f.add_paragraph()
    p2.text = "similarity = 0.5 × track alignment + 0.5 × vector semantic overlap (with a small beginner-friendliness bonus)"
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_DARK
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    score_cards = [
        ("Track Similarity (50%)", "Matches goal language (\"ML engineer\", \"data scientist\", \"doctor\", \"lawyer\") to each course's canonical track."),
        ("Vector Cosine Overlap (50%)", "Compares goal terms against course titles and descriptions using 384-dimensional vector embeddings in pgvector."),
        ("Skill-Gap Score (40%)", "Rewards courses that teach competencies the learner doesn't yet have, accelerating tangible career readiness.")
    ]

    for idx, (stitle, sdesc) in enumerate(score_cards):
        s_left = Inches(0.8) + idx * (col_w + gap)
        add_card(s7, s_left, Inches(3.35), col_w, Inches(3.4))
        tf = create_textbox(s7, s_left + Inches(0.3), Inches(3.6), col_w - Inches(0.6), Inches(2.8))
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(14.5)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(100, 100, 95)
        p2.space_before = Pt(12)

    # ==========================================
    # SLIDE 8: Multidisciplinary (Dark)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8, DARK_BG)
    add_header(s8, "Multi-Disciplinary Taxonomy", "Broadening beyond tech: High-demand CS & Non-CS tracks", dark=True)

    tracks = [
        ("💻 Technology & AI", ["Machine Learning & Deep Learning", "Full-Stack Web (React/Node)", "Cloud & DevOps (Docker/K8s)", "Data Science & Python"]),
        ("🩺 Healthcare & Medicine", ["Human Anatomy & Physiology", "Clinical Surgery Fundamentals", "Medical Ethics & Bioethics", "Patient Health Assessment"]),
        ("⚖️ Law & Legal Ops", ["Criminal Law & Litigation", "Corporate Contract Law", "Legal Negotiation & ADR", "Constitutional Jurisprudence"]),
        ("🏗️ Core Engineering", ["Civil Structural Design", "Mechanical Thermal Systems", "AutoCAD & Engineering Drafting", "Material Science Fundamentals"]),
        ("🔬 Forensics & Psychology", ["Crime Scene Forensics & Evidence", "Cognitive & Clinical Psychology", "DNA Profiling & Ballistics", "Behavioral & Criminal Analysis"]),
        ("📈 Creator & Digital Economy", ["Digital Marketing Strategy", "Growth Analytics & SEO", "Content Creation & Storytelling", "Brand Strategy & Monetization"])
    ]

    for idx, (tname, tskills) in enumerate(tracks):
        row = idx // 3
        col = idx % 3
        t_left = Inches(0.8) + col * (card_w5 + Inches(0.38))
        t_top = Inches(1.9) + row * (card_h5 + Inches(0.3))

        add_card(s8, t_left, t_top, card_w5, card_h5, bg_color=DARK_CARD_BG, border_color=PRIMARY)
        tf = create_textbox(s8, t_left + Inches(0.25), t_top + Inches(0.18), card_w5 - Inches(0.5), card_h5 - Inches(0.36))
        p = tf.paragraphs[0]
        p.text = tname
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = SECONDARY
        for sk in tskills:
            p_sk = tf.add_paragraph()
            p_sk.text = f"• {sk}"
            p_sk.font.size = Pt(10.5)
            p_sk.font.color.rgb = RGBColor(210, 215, 205)

    # ==========================================
    # SLIDE 9: Data Model (Light)
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9, LIGHT_BG)
    add_header(s9, "Data Model", "Nine tables, RLS on everything personal")

    t_shape = s9.shapes.add_table(10, 3, Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.9))
    tbl = t_shape.table
    tbl.columns[0].width = Inches(2.6)
    tbl.columns[1].width = Inches(7.13)
    tbl.columns[2].width = Inches(2.0)

    headers = ["Table", "Purpose", "Access"]
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        tf_c = cell.text_frame
        tf_c.margin_left = Inches(0.1)
        p = tf_c.paragraphs[0]
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT

    rows_data = [
        ("learner_profiles", "Goal, experience level, preferences, & interests.", "RLS"),
        ("skills", "Canonical skill catalogue across CS & Non-CS domains.", "Public read"),
        ("learner_skills", "Per-user proficiency (self-reported / inferred / assessed).", "RLS"),
        ("courses", "Catalog with pgvector(384) embeddings.", "Public read"),
        ("course_skills", "Course ↔ skill mapping including prerequisites.", "Public read"),
        ("learning_paths", "Ordered, milestone-labeled curricula per learner.", "RLS"),
        ("path_items", "Ordered milestone courses, statuses, scores, & AI explanations.", "RLS"),
        ("quiz_attempts", "Server-side questions, user submissions, & verification grades.", "RLS"),
        ("chat_history", "AI guide conversation log per learner.", "RLS")
    ]

    for row_idx, r in enumerate(rows_data):
        for col_idx, val in enumerate(r):
            cell = tbl.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 == 0 else RGBColor(245, 243, 238)
            tf_c = cell.text_frame
            tf_c.margin_left = Inches(0.1)
            p = tf_c.paragraphs[0]
            p.text = val
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_DARK if col_idx < 2 else (PRIMARY if val == "RLS" else SECONDARY)
            p.font.bold = (col_idx == 0 or col_idx == 2)

    # ==========================================
    # SLIDE 10: User Journey (Dark)
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10, DARK_BG)
    add_header(s10, "User Journey", "From sign-in to a living, adaptive curriculum", dark=True)

    journey_steps = [
        ("1", "Login", "Supabase Auth sign-in / sign-up with session security."),
        ("2", "Onboarding / Resume", "Learner states their goal or uploads resume (.pdf / .docx)."),
        ("3", "Dashboard", "Mastery radar, active paths, skill gaps, & progress bars."),
        ("4", "Path Detail", "Milestone timeline, course rationales, & Edit Plan pruning."),
        ("5", "Verify & Adapt", "Interactive MCQ quizzes, skill boosting, & AI guide chat.")
    ]

    j_w = Inches(2.15)
    j_gap = Inches(0.24)
    j_top = Inches(2.0)
    j_h = Inches(4.5)

    for i, (num, j_title, j_desc) in enumerate(journey_steps):
        left_pos = Inches(0.8) + i * (j_w + j_gap)
        add_card(s10, left_pos, j_top, j_w, j_h, bg_color=DARK_CARD_BG, border_color=PRIMARY)

        b_num = s10.shapes.add_shape(MSO_SHAPE.OVAL, left_pos + Inches(0.2), j_top + Inches(0.35), Inches(0.55), Inches(0.55))
        b_num.fill.solid()
        b_num.fill.fore_color.rgb = SECONDARY
        b_num.line.fill.background()
        b_num.text = num
        b_num.text_frame.paragraphs[0].font.size = Pt(13)
        b_num.text_frame.paragraphs[0].font.bold = True
        b_num.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
        b_num.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        tf = create_textbox(s10, left_pos + Inches(0.2), j_top + Inches(1.1), j_w - Inches(0.4), Inches(3.0))
        p1 = tf.paragraphs[0]
        p1.text = j_title
        p1.font.size = Pt(14.5)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_LIGHT

        p2 = tf.add_paragraph()
        p2.text = j_desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(200, 205, 195)
        p2.space_before = Pt(8)

    tf_j_foot = create_textbox(s10, Inches(0.8), Inches(6.65), Inches(11.73), Inches(0.35))
    p = tf_j_foot.paragraphs[0]
    p.text = "Every stage is protected by session-aware routing; unauthenticated visitors are redirected straight back to Login."
    p.font.size = Pt(10.5)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 11: Conclusion (Dark)
    # ==========================================
    s11 = prs.slides.add_slide(blank_layout)
    set_bg(s11, DARK_BG)

    badge11 = s11.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.16), Inches(1.6), Inches(1.0), Inches(1.0))
    badge11.fill.solid()
    badge11.fill.fore_color.rgb = SECONDARY
    badge11.line.fill.background()
    badge11.text = "🌱"
    badge11.text_frame.paragraphs[0].font.size = Pt(28)
    badge11.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    tf11 = create_textbox(s11, Inches(1.0), Inches(2.9), Inches(11.33), Inches(1.2))
    p = tf11.paragraphs[0]
    p.text = "Grow a curriculum around every learner"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    tf_s11 = create_textbox(s11, Inches(1.5), Inches(4.2), Inches(10.33), Inches(1.0))
    p = tf_s11.paragraphs[0]
    p.text = "LearnAI — Conversational goals, resume skill-gap intelligence, explainable recommendations, and a design system built to feel human."
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(200, 205, 195)
    p.alignment = PP_ALIGN.CENTER

    tf_ty = create_textbox(s11, Inches(1.0), Inches(5.6), Inches(11.33), Inches(0.8))
    p = tf_ty.paragraphs[0]
    p.text = "Thank you — Code Catalyst"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    out_file = "LearnAI_Deck_v2.pptx"
    prs.save(out_file)
    print(f"Updated presentation saved to: {os.path.abspath(out_file)}")
    try:
        prs.save("LearnAI_Presentation.pptx")
        print("Also updated LearnAI_Presentation.pptx")
    except Exception:
        print("Note: LearnAI_Presentation.pptx is currently open in another app.")

if __name__ == "__main__":
    create_presentation()
